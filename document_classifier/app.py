from flask import Flask, request, jsonify, render_template_string
import os
import re
import json
import csv
import io
import pandas as pd
import fitz  # PyMuPDF
import docx
import pptx
from bs4 import BeautifulSoup
import xml.etree.ElementTree as ET
from striprtf.striprtf import rtf_to_text

app = Flask(__name__)

PORT = 5003

# Define regex patterns for PII detection
PII_PATTERNS = {
    "email": r"[a-zA-Z0-9_.+-]+@[a-zA-Z0-9-]+\.[a-zA-Z0-9-.]+",
    "phone": r"\+?\d[\d\s\-]{8,15}",
    "aadhaar": r"\b\d{4}[\s-]?\d{4}[\s-]?\d{4}\b",
    "pan": r"\b[A-Z]{5}[0-9]{4}[A-Z]\b",
    "passport": r"\b[A-Z]{1}-?\d{7}\b",
    "ssn": r"\b\d{3}-\d{2}-\d{4}\b",
    "ifsc": r"\b[A-Z]{4}0[A-Z0-9]{6}\b",
    "credit_card": r"\b(?:\d[ -]*?){13,16}\b",
    "ip_address": r"\b(?:\d{1,3}\.){3}\d{1,3}\b",
    "mac_address": r"\b(?:[0-9A-Fa-f]{2}[:-]){5}[0-9A-Fa-f]{2}\b",
    "dob": r"\b(?:\d{1,2}[-/th|st|nd|rd\s]*)?(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*[-/\s]*\d{2,4}\b|\b\d{1,2}[-/]\d{1,2}[-/]\d{2,4}\b",
    "gender": r"\b(?:male|female|other|non-binary|transgender)\b",
    "name": r"\b[A-Z][a-z]+\s[A-Z][a-z]+\b",
    "address": r"\d{1,5}\s\w+\s\w+",
    "voter_id": r"\b[A-Z]{3}[0-9]{7}\b",
    "bank_account": r"\b\d{9,18}\b",
    "vehicle_reg": r"\b[A-Z]{2}[0-9]{2}[A-Z]{2}[0-9]{4}\b",
    "employee_id": r"\bEMP[0-9]{4,6}\b",
    "medical_record": r"\bMRN[0-9]{6,8}\b",
    "insurance_policy": r"\b[A-Z]{2}[0-9]{10}\b",
}


# Extract text from different file types
def extract_text(file_stream, filename):
    ext = os.path.splitext(filename)[1].lower()
    if ext == ".pdf":
        doc = fitz.open(stream=file_stream.read(), filetype="pdf")
        return "\n".join([page.get_text() for page in doc])
    elif ext == ".docx":
        doc = docx.Document(file_stream)
        return "\n".join([para.text for para in doc.paragraphs])
    elif ext == ".pptx":
        prs = pptx.Presentation(file_stream)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    elif ext == ".txt":
        return file_stream.read().decode("utf-8", errors="ignore")
    elif ext == ".rtf":
        return rtf_to_text(file_stream.read().decode("utf-8", errors="ignore"))
    elif ext in [".csv", ".tsv"]:
        delimiter = "," if ext == ".csv" else "\t"
        decoded = file_stream.read().decode("utf-8", errors="ignore")
        reader = csv.reader(io.StringIO(decoded), delimiter=delimiter)
        return "\n".join(["\t".join(row) for row in reader])
    elif ext == ".xlsx":
        df = pd.read_excel(file_stream, engine="openpyxl")
        return df.to_csv(index=False, sep="\t")
    elif ext == ".json":
        data = json.load(file_stream)
        return json.dumps(data, indent=2)
    elif ext == ".xml":
        tree = ET.parse(file_stream)
        root = tree.getroot()
        return ET.tostring(root, encoding="unicode")
    elif ext == ".html":
        soup = BeautifulSoup(file_stream.read(), "html.parser")
        return soup.get_text()
    else:
        return ""


# Detect PII using regex
def detect_pii(text, selected_types=None):
    lines = text.splitlines()
    pii_found = {}
    locations = {}
    selected_patterns = {
        k: v
        for k, v in PII_PATTERNS.items()
        if (selected_types is None or k in selected_types)
    }

    for i, line in enumerate(lines, start=1):
        for pii_type, pattern in selected_patterns.items():
            matches = re.findall(pattern, line, flags=re.IGNORECASE)
            if matches:
                pii_found.setdefault(pii_type, []).extend(matches)
                locations.setdefault(pii_type, []).append(f"Line {i}")
    return pii_found, locations


@app.route("/document-upload", methods=["POST"])
def document_upload():
    uploaded_files = request.files.getlist("files")
    selected_pii_types = request.form.getlist("pii_types")

    if not selected_pii_types:
        selected_pii_types = list(PII_PATTERNS.keys())

    if not uploaded_files:
        return jsonify({"error": "No files uploaded"}), 400

    results = []
    for uploaded_file in uploaded_files:
        try:
            file_stream = io.BytesIO(uploaded_file.read())
            file_stream.seek(0)
            text = extract_text(file_stream, uploaded_file.filename)
            pii_types, locations = detect_pii(text, selected_types=selected_pii_types)
            summary = {
                pii_type: len(matches) for pii_type, matches in pii_types.items()
            }
            results.append(
                {
                    "file_name": uploaded_file.filename,
                    "pii_found": bool(pii_types),
                    "classifications": summary,
                }
            )

        except Exception as e:
            results.append({"file_name": uploaded_file.filename, "error": str(e)})

    return jsonify(results)


if __name__ == "__main__":
    app.run(host="0.0.0.0", port=5003, debug=True)

    # Homepage with file upload and PII selection

# @app.route('/')
# def home():
#     return render_template_string("""
#         <!DOCTYPE html>
#         <html>
#         <head><title>PII Detector</title></head>
#         <body>
#             <h2>Upload Documents to Detect PII</h2>
#             <form method="POST" action="/document-upload" enctype="multipart/form-data">
#                 <input type="file" name="files" multiple required><br><br>
#                 <label><input type="checkbox" name="pii_types" value="aadhaar"> Aadhaar</label><br>
#                 <label><input type="checkbox" name="pii_types" value="pan"> PAN</label><br>
#                 <label><input type="checkbox" name="pii_types" value="passport"> Passport</label><br>
#                 <label><input type="checkbox" name="pii_types" value="email"> Email</label><br>
#                 <label><input type="checkbox" name="pii_types" value="phone"> Phone</label><br>
#                 <br><button type="submit">Upload</button>
#             </form>
#         </body>
#         </html>
#     """)
